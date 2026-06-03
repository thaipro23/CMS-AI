from __future__ import annotations

import base64
import binascii
import csv
import json
import re
import zipfile
from dataclasses import dataclass
from html import unescape
from io import BytesIO, StringIO
from pathlib import PurePosixPath
import xml.etree.ElementTree as ET
from typing import Any

from bs4 import BeautifulSoup

from app.core.config import settings
from app.services.problem_parser import build_ai_text_from_problem, remove_openedx_filename_metadata
from pypdf import PdfReader
from pptx import Presentation


try:  # Optional OCR/rendering stack for v25.9.13.0.
    from PIL import Image, ImageOps
except Exception:  # pragma: no cover - optional dependency guard
    Image = None
    ImageOps = None

try:  # Optional OCR engine wrapper. Requires tesseract binary in the image/host.
    import pytesseract
except Exception:  # pragma: no cover - optional dependency guard
    pytesseract = None

try:  # Fast PDF renderer when available.
    import fitz  # PyMuPDF
except Exception:  # pragma: no cover - optional dependency guard
    fitz = None

try:  # Fallback PDF renderer; requires poppler-utils.
    from pdf2image import convert_from_bytes
except Exception:  # pragma: no cover - optional dependency guard
    convert_from_bytes = None

try:  # Optional but installed in the v25.9.12.8 backend requirements.
    from docx import Document
except Exception:  # pragma: no cover - optional dependency guard for old containers
    Document = None

try:  # Optional but installed in the v25.9.12.8 backend requirements.
    from openpyxl import load_workbook
except Exception:  # pragma: no cover - optional dependency guard for old containers
    load_workbook = None


@dataclass
class ExtractedContent:
    block_id: str
    block_type: str
    display_name: str
    content: str
    source_type: str
    source_ref: str
    page_number: int | None = None
    timestamp_start: str | None = None
    timestamp_end: str | None = None
    parent_block_id: str | None = None


class ContentExtractor:
    """Normalize Open edX blocks/assets into plain text with source references.

    v25.9.12.8 adds broad handout/upload extraction.  It supports CMS-linked
    assets and teacher-uploaded node files in these families:

    - PDF: .pdf; v25.9.13.0 can OCR scanned pages when FILE_OCR_ENABLED=true
    - PowerPoint: .pptx/.ppt when python-pptx can open the file; v25.9.13.0 also reads speaker notes and can OCR slide images when enabled
    - Word: .docx
    - Excel: .xlsx/.xlsm
    - CSV/TSV: .csv/.tsv
    - Text/Markdown/HTML/XML/JSON/WebVTT/SRT: text-based parsing

    Legacy binary .doc/.xls files are detected but cannot be parsed reliably in a
    pure Python/offline container.  Uploaders should convert them to .docx/.xlsx
    or PDF.  The error is returned explicitly by the upload API.
    """

    TEXT_BLOCK_TYPES = {"html", "text", "problem", "vertical", "sequential", "chapter"}
    VIDEO_BLOCK_TYPES = {"video", "transcript"}
    FILE_BLOCK_TYPES = {
        "pdf", "ppt", "pptx", "doc", "docx", "xls", "xlsx", "xlsm", "csv", "tsv",
        "txt", "md", "markdown", "html", "htm", "json", "xml", "vtt", "srt",
        "file", "asset", "handout", "document",
    }

    def extract_blocks(self, blocks: list[dict[str, Any]]) -> list[ExtractedContent]:
        extracted: list[ExtractedContent] = []
        for block in blocks:
            extracted.extend(self.extract_block(block))
        return [item for item in extracted if item.content.strip()]

    def extract_block(self, block: dict[str, Any]) -> list[ExtractedContent]:
        block_id = str(block.get("block_id") or block.get("id") or "")
        if not block_id:
            return []
        block_type = str(block.get("type") or block.get("block_type") or "unknown").lower()
        display_name = str(block.get("display_name") or block.get("title") or block_type)
        parent = block.get("parent_block_id") or block.get("parent")
        source_ref = str(block.get("source_ref") or block_id)

        candidates = [
            # Studio connector can expose raw problem XML in a dedicated field.
            # Prefer it so we can preserve correct="true" before converting to text.
            block.get("problem_xml"),
            block.get("data"),
            block.get("content"),
            block.get("html"),
            block.get("transcript"),
            block.get("student_view_data"),
        ]
        raw_content = remove_openedx_filename_metadata(self._best_text(candidates))
        if block_type == "problem":
            content = build_ai_text_from_problem(raw_content) or self.clean_html(raw_content)
        elif block_type in self.TEXT_BLOCK_TYPES:
            content = self.clean_html(raw_content)
        elif block_type in self.VIDEO_BLOCK_TYPES:
            content = self.clean_transcript(raw_content)
        else:
            content = self.clean_html(raw_content)
        content = remove_openedx_filename_metadata(content)

        result = []
        if content:
            result.append(ExtractedContent(
                block_id=block_id,
                block_type=block_type,
                display_name=display_name,
                content=content,
                source_type=self._source_type(block_type),
                source_ref=source_ref,
                timestamp_start=block.get("timestamp_start"),
                timestamp_end=block.get("timestamp_end"),
                parent_block_id=str(parent) if parent else None,
            ))

        # Connector plugin can attach expanded assets/transcripts to any block.
        for asset in block.get("assets") or []:
            result.extend(self.extract_asset(asset, parent_block_id=block_id))
        for transcript in block.get("transcripts") or []:
            result.extend(self.extract_block({**transcript, "type": "transcript", "parent_block_id": block_id}))
        return result

    def extract_asset(self, asset: dict[str, Any], parent_block_id: str | None = None) -> list[ExtractedContent]:
        asset_id = str(asset.get("asset_id") or asset.get("id") or asset.get("url") or asset.get("source_ref") or "asset")
        filename = str(asset.get("file_name") or asset.get("filename") or asset.get("display_name") or self._filename_from_ref(asset_id) or asset_id)
        mime = str(asset.get("mime_type") or asset.get("content_type") or "").lower()
        raw_bytes = self._asset_bytes(asset)
        raw_text = asset.get("text") or asset.get("content") or ""
        ext = self._extension(filename, mime, asset_id)

        pages: list[tuple[int | None, str]] = []
        if raw_bytes:
            try:
                pages = self.extract_file_pages(raw_bytes, filename=filename, mime_type=mime, ext=ext)
            except ValueError:
                if asset.get('strict'):
                    raise
                pages = []
            except Exception:
                if asset.get('strict'):
                    raise
                pages = []
        if not pages and raw_text:
            text = self.clean_html(str(raw_text))
            if text:
                pages = [(None, text)]

        items: list[ExtractedContent] = []
        for page, text in pages:
            if not text.strip():
                continue
            page_block_id = f"{asset_id}#page={page}" if page else asset_id
            source_ref = str(asset.get("source_ref") or asset.get("url") or asset_id)
            if page:
                source_ref = f"{source_ref}#page={page}"
            items.append(ExtractedContent(
                block_id=page_block_id,
                block_type=ext or "file",
                display_name=f"{filename} - page {page}" if page else filename,
                content=text,
                source_type=ext or "file",
                source_ref=source_ref,
                page_number=page,
                parent_block_id=parent_block_id,
            ))
        return items

    def extract_file_pages(self, raw: bytes, *, filename: str = '', mime_type: str = '', ext: str = '') -> list[tuple[int | None, str]]:
        ext = (ext or self._extension(filename, mime_type)).lower()
        mime = (mime_type or '').lower()
        if ext == "pdf" or "pdf" in mime:
            return self.extract_pdf_pages(raw)
        if ext in {"pptx", "ppt"} or "presentation" in mime:
            return self.extract_pptx_pages(raw)
        if ext == "docx" or "wordprocessingml" in mime:
            text = self.extract_docx_text(raw)
            return [(None, text)] if text else []
        if ext in {"xlsx", "xlsm"} or "spreadsheetml" in mime:
            return self.extract_xlsx_sheets(raw)
        if ext in {"csv", "tsv"} or "csv" in mime or "tab-separated" in mime:
            text = self.extract_csv_text(raw, delimiter='\t' if ext == 'tsv' else None)
            return [(None, text)] if text else []
        if ext in {"txt", "md", "markdown", "html", "htm", "xml", "json", "srt", "vtt"} or mime.startswith('text/') or 'json' in mime or 'xml' in mime:
            text = self.extract_text_file(raw, ext=ext, mime_type=mime)
            return [(None, text)] if text else []
        if ext in {"doc", "xls"}:
            raise ValueError(f'File {filename or ext} là định dạng Office cũ .{ext}. Hãy đổi sang .docx/.xlsx hoặc PDF để hệ thống tách nội dung ổn định.')
        # Last resort: try text decoding for unknown but text-like uploads.
        text = self.extract_text_file(raw, ext=ext, mime_type=mime)
        return [(None, text)] if text else []

    def extract_pdf_pages(self, raw: bytes) -> list[tuple[int, str]]:
        reader = PdfReader(BytesIO(raw))
        page_texts: dict[int, str] = {}
        empty_pages: list[int] = []
        for index, page in enumerate(reader.pages, start=1):
            text = self.normalize_text(page.extract_text() or "")
            if text:
                page_texts[index] = text
            else:
                empty_pages.append(index)

        if empty_pages and settings.file_ocr_enabled:
            for page_number, text in self._ocr_pdf_pages(raw, empty_pages).items():
                if text:
                    page_texts[page_number] = text

        return [(page_number, page_texts[page_number]) for page_number in sorted(page_texts)]

    def extract_pptx_pages(self, raw: bytes) -> list[tuple[int, str]]:
        pres = Presentation(BytesIO(raw))
        notes_by_slide = self._extract_pptx_notes(raw) if settings.pptx_extract_speaker_notes else {}
        slides = []
        for index, slide in enumerate(pres.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text)
                image_text = self._ocr_pptx_shape_image(shape)
                if image_text:
                    parts.append(f"OCR ảnh trong slide: {image_text}")
            notes = notes_by_slide.get(index)
            if notes:
                parts.append(f"Ghi chú diễn giả: {notes}")
            text = self.normalize_text("\n".join(parts))
            if text:
                slides.append((index, text))
        return slides


    def _ocr_pdf_pages(self, raw: bytes, page_numbers: list[int]) -> dict[int, str]:
        if not settings.file_ocr_enabled or pytesseract is None or Image is None:
            return {}
        max_pages = max(0, int(settings.file_ocr_max_pages or 0))
        selected_pages = page_numbers[:max_pages] if max_pages else []
        if not selected_pages:
            return {}

        texts: dict[int, str] = {}
        if fitz is not None:
            try:
                document = fitz.open(stream=raw, filetype='pdf')
                for page_number in selected_pages:
                    page = document.load_page(page_number - 1)
                    pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
                    image = Image.open(BytesIO(pixmap.tobytes('png')))
                    text = self._ocr_image(image)
                    if text:
                        texts[page_number] = text
                document.close()
                return texts
            except Exception:
                texts = {}

        if convert_from_bytes is not None:
            for page_number in selected_pages:
                try:
                    images = convert_from_bytes(raw, dpi=200, first_page=page_number, last_page=page_number)
                    if images:
                        text = self._ocr_image(images[0])
                        if text:
                            texts[page_number] = text
                except Exception:
                    continue
        return texts

    def _ocr_pptx_shape_image(self, shape: Any) -> str:
        if not settings.file_ocr_enabled or not settings.pptx_ocr_images_enabled:
            return ''
        if pytesseract is None or Image is None or not hasattr(shape, 'image'):
            return ''
        try:
            return self._ocr_image(Image.open(BytesIO(shape.image.blob)))
        except Exception:
            return ''

    def _ocr_image(self, image: Any) -> str:
        if pytesseract is None:
            return ''
        try:
            prepared = self._prepare_ocr_image(image)
            configs = [settings.file_ocr_tesseract_config or '--oem 3 --psm 6']
            # Fallback modes: useful for sparse diagrams/screenshots where a single PSM can miss text.
            for fallback in ('--oem 3 --psm 6', '--oem 3 --psm 11'):
                if fallback not in configs:
                    configs.append(fallback)

            candidates = []
            for config in configs:
                try:
                    text = pytesseract.image_to_string(
                        prepared,
                        lang=settings.file_ocr_language or 'vie+eng',
                        config=config,
                    )
                    normalized = self.normalize_text(text)
                    if normalized:
                        candidates.append(normalized)
                except Exception:
                    continue
            return max(candidates, key=len) if candidates else ''
        except Exception:
            return ''

    def _prepare_ocr_image(self, image: Any) -> Any:
        if Image is None or ImageOps is None:
            return image
        try:
            prepared = image.convert('RGB') if getattr(image, 'mode', '') not in {'RGB', 'L'} else image
            # Grayscale improves OCR stability for rendered PDF screenshots and dark-theme slides.
            prepared = ImageOps.grayscale(prepared)
            return prepared
        except Exception:
            return image

    def _extract_pptx_notes(self, raw: bytes) -> dict[int, str]:
        """Extract speaker notes from ppt/notesSlides using slide relationships."""
        try:
            archive = zipfile.ZipFile(BytesIO(raw))
        except Exception:
            return {}

        slide_files = sorted(
            [name for name in archive.namelist() if re.match(r'ppt/slides/slide\d+\.xml$', name)],
            key=lambda value: int(re.search(r'slide(\d+)\.xml$', value).group(1)),
        )
        notes: dict[int, str] = {}
        for slide_index, slide_path in enumerate(slide_files, start=1):
            rels_path = f"ppt/slides/_rels/{PurePosixPath(slide_path).name}.rels"
            if rels_path not in archive.namelist():
                continue
            try:
                rel_root = ET.fromstring(archive.read(rels_path))
            except Exception:
                continue
            note_path = ''
            for rel in rel_root:
                rel_type = rel.attrib.get('Type', '')
                target = rel.attrib.get('Target', '')
                if 'notesSlide' in rel_type and target:
                    note_path = str((PurePosixPath(slide_path).parent / target).as_posix())
                    break
            if not note_path:
                continue
            # Normalize paths such as ppt/slides/../notesSlides/notesSlide1.xml.
            parts: list[str] = []
            for part in PurePosixPath(note_path).parts:
                if part == '..' and parts:
                    parts.pop()
                elif part not in {'.', ''}:
                    parts.append(part)
            note_path = '/'.join(parts)
            if note_path not in archive.namelist():
                continue
            text = self._xml_text(archive.read(note_path))
            if text:
                notes[slide_index] = text
        archive.close()
        return notes

    def _xml_text(self, raw: bytes) -> str:
        try:
            root = ET.fromstring(raw)
        except Exception:
            return ''
        texts: list[str] = []
        for node in root.iter():
            if node.tag.endswith('}t') or node.tag == 't':
                if node.text and node.text.strip():
                    texts.append(node.text.strip())
        return self.normalize_text('\n'.join(texts))

    def extract_docx_text(self, raw: bytes) -> str:
        if Document is None:
            raise ValueError('Thiếu package python-docx trong backend. Cài lại image bằng docker compose up --build.')
        doc = Document(BytesIO(raw))
        parts: list[str] = []
        for paragraph in doc.paragraphs:
            text = paragraph.text.strip()
            if text:
                parts.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
                if cells:
                    parts.append(' | '.join(cells))
        return self.normalize_text('\n'.join(parts))

    def extract_xlsx_sheets(self, raw: bytes) -> list[tuple[int | None, str]]:
        if load_workbook is None:
            raise ValueError('Thiếu package openpyxl trong backend. Cài lại image bằng docker compose up --build.')
        workbook = load_workbook(BytesIO(raw), read_only=True, data_only=True)
        sheets: list[tuple[int | None, str]] = []
        for index, sheet in enumerate(workbook.worksheets, start=1):
            lines: list[str] = [f'Sheet: {sheet.title}']
            for row in sheet.iter_rows(values_only=True):
                cells = [self._cell_to_text(cell) for cell in row]
                # trim trailing empty cells
                while cells and not cells[-1]:
                    cells.pop()
                if any(cells):
                    lines.append(' | '.join(cells))
            text = self.normalize_text('\n'.join(lines))
            if text and text != f'Sheet: {sheet.title}':
                sheets.append((index, text))
        return sheets

    def extract_csv_text(self, raw: bytes, delimiter: str | None = None) -> str:
        text = self._decode_text_bytes(raw)
        sample = text[:4096]
        if delimiter is None:
            try:
                dialect = csv.Sniffer().sniff(sample)
                delimiter = dialect.delimiter
            except Exception:
                delimiter = ','
        reader = csv.reader(StringIO(text), delimiter=delimiter)
        lines: list[str] = []
        for row in reader:
            cells = [self.normalize_text(str(cell)) for cell in row]
            if any(cells):
                lines.append(' | '.join(cells))
        return self.normalize_text('\n'.join(lines))

    def extract_text_file(self, raw: bytes, *, ext: str = '', mime_type: str = '') -> str:
        text = self._decode_text_bytes(raw)
        if ext in {'html', 'htm'} or 'html' in mime_type:
            return self.clean_html(text)
        if ext == 'json' or 'json' in mime_type:
            try:
                parsed = json.loads(text)
                text = json.dumps(parsed, ensure_ascii=False, indent=2)
            except Exception:
                pass
            return self.normalize_text(text)
        if ext in {'xml'} or 'xml' in mime_type:
            return self.clean_html(text)
        if ext in {'srt', 'vtt'}:
            return self.clean_transcript(text)
        return self.normalize_text(text)

    def clean_html(self, value: Any) -> str:
        text = remove_openedx_filename_metadata(self._best_text([value]))
        soup = BeautifulSoup(text, "html.parser")
        for tag in soup(["script", "style"]):
            tag.extract()
        return remove_openedx_filename_metadata(self.normalize_text(soup.get_text(" ")))

    def clean_transcript(self, value: Any) -> str:
        text = self._best_text([value])
        text = re.sub(r'^WEBVTT.*?\n', ' ', text, flags=re.IGNORECASE | re.DOTALL)
        text = re.sub(r"\d+\s*\n", " ", text)
        text = re.sub(r"\d{1,2}:\d{2}:\d{2}[,.]\d{1,3}\s*-->\s*\d{1,2}:\d{2}:\d{2}[,.]\d{1,3}[^\n]*", " ", text)
        return self.normalize_text(text)

    def normalize_text(self, text: str) -> str:
        text = remove_openedx_filename_metadata(unescape(text or ""))
        text = re.sub(r"\s+", " ", text)
        return remove_openedx_filename_metadata(text.strip())

    def _best_text(self, values: list[Any]) -> str:
        for value in values:
            if value is None:
                continue
            if callable(value):
                continue
            if isinstance(value, dict):
                nested = self._best_text([value.get("text"), value.get("html"), value.get("content"), value.get("transcript")])
                if nested:
                    return nested
                continue
            if isinstance(value, list):
                joined = "\n".join(self._best_text([item]) for item in value)
                if joined.strip():
                    return joined
                continue
            text = str(value)
            if text.strip():
                return text
        return ""

    def _source_type(self, block_type: str) -> str:
        if block_type in {"html", "text"}:
            return "html"
        if block_type in {"video", "transcript"}:
            return "transcript"
        if block_type == "problem":
            return "problem"
        if block_type in self.FILE_BLOCK_TYPES:
            return block_type if block_type not in {"asset", "handout", "document"} else "file"
        return block_type or "unknown"

    def _asset_bytes(self, asset: dict[str, Any]) -> bytes | None:
        raw = asset.get('bytes') or asset.get('raw_bytes')
        if isinstance(raw, bytes):
            return raw
        if isinstance(raw, bytearray):
            return bytes(raw)
        if isinstance(raw, str) and raw.strip():
            try:
                return base64.b64decode(raw, validate=True)
            except Exception:
                return raw.encode('utf-8', errors='ignore')
        encoded = asset.get('bytes_base64') or asset.get('content_base64') or asset.get('data_base64')
        if isinstance(encoded, str) and encoded.strip():
            try:
                return base64.b64decode(encoded, validate=True)
            except (binascii.Error, ValueError):
                return None
        return None

    def _decode_text_bytes(self, raw: bytes) -> str:
        for encoding in ('utf-8-sig', 'utf-8', 'utf-16', 'cp1258', 'latin-1'):
            try:
                return raw.decode(encoding)
            except Exception:
                continue
        return raw.decode('utf-8', errors='ignore')

    @staticmethod
    def _cell_to_text(value: Any) -> str:
        if value is None:
            return ''
        return str(value).strip()

    def _extension(self, filename: str = '', mime: str = '', ref: str = '') -> str:
        source = f'{filename} {ref}'.lower()
        # Extract from normal filename or Open edX asset locator block@file.ext.
        match = re.search(r'\.([a-z0-9]{1,8})(?:[?#\s]|$)', source)
        if match:
            return match.group(1)
        if 'pdf' in mime:
            return 'pdf'
        if 'wordprocessingml' in mime:
            return 'docx'
        if 'spreadsheetml' in mime:
            return 'xlsx'
        if 'presentationml' in mime or 'powerpoint' in mime:
            return 'pptx'
        if 'csv' in mime:
            return 'csv'
        if 'html' in mime:
            return 'html'
        if 'json' in mime:
            return 'json'
        if 'xml' in mime:
            return 'xml'
        if mime.startswith('text/'):
            return 'txt'
        return ''

    @staticmethod
    def _filename_from_ref(ref: str) -> str:
        if not ref:
            return ''
        cleaned = ref.split('?', 1)[0].split('#', 1)[0].rstrip('/')
        if 'block@' in cleaned:
            return cleaned.rsplit('block@', 1)[-1]
        return cleaned.rsplit('/', 1)[-1]
